from PyQt5.QtCore import QThread, pyqtSignal
import pandas as pd
from pptx import Presentation
from datetime import datetime
from comtypes import client
import os
import re
import logging

import win32com.client as win32


logging.basicConfig(
    filename='certificates_generation.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

class CertificateWorker(QThread):
    progress = pyqtSignal(str)
    finished = pyqtSignal(str)
    progress_percent = pyqtSignal(int)   # progression en pourcentage


    #def __init__(self, template_path, excel_path, formation_title, output_folder):
    def __init__(self, template_path, excel_path, formation_title, output_folder,score_min=80, date_start=None, date_end=None):
        super().__init__()
        self.template_path = template_path
        self.excel_path = excel_path
        self.formation_title = formation_title
        self.output_folder = output_folder
# Initialiser 
        self.score_min = score_min
        self.date_start = date_start
        self.date_end = date_end
        self.total_certificates = 0

    def clean_filename(self, s):
        return re.sub(r'[\\/*?:"<>|]', "", s)

    def format_date_with_suffix(self, date_obj):
        day = date_obj.day
        suffix = "th" if 11 <= day <= 13 else {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
        return f"{day}{suffix} {date_obj.strftime('%B %Y')}"

    def replace_text(self, prs, replacements):
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, val in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, val)

    def convert_to_pdf(self, pptx_path, pdf_path):
        powerpoint = client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # Affiche PowerPoint (utile pour le debug)
        try:
            presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=True)
            presentation.SaveAs(pdf_path, 32)  # 32 = format PDF
            presentation.Close()
            #┘ self.progress.emit(f" PDF created : {pdf_path}")
        except Exception as e:
            self.progress.emit(f" PDF conversion error: {e}")
        finally:
            powerpoint.Quit()

    def envoyer_email_outlook(self, destinataire, sujet, corps, fichier_joint=None, cc=None):
        try:
            import re
            import getpass

            # 1. Récupérer le dossier des signatures Outlook
            signature_dir = os.path.join(os.environ['APPDATA'], "Microsoft", "Signatures")
            username = getpass.getuser()  # Récupère le nom d'utilisateur Windows

            # 2. Lister toutes les signatures disponibles
            signature_files = os.listdir(signature_dir)
            signature_path = None

            # 3. Chercher une signature contenant le nom d'utilisateur
            for signature_file in signature_files:
                if username in signature_file:  # Si le nom d'utilisateur est dans le nom du fichier
                    signature_path = os.path.join(signature_dir, signature_file)
                    break  # Sortir dès qu'une signature est trouvée

            if signature_path is None:
                logging.error(f"No signature found for username '{username}.")
            else:
                try:
                    # 4. Lire la signature HTML
                    with open(signature_path, "r", encoding="mbcs") as f:
                        signature_html = f.read()

                    # 5. Chercher les images référencées dans la signature
                    matches = re.findall(r'<img[^>]+src="([^"]+)"', signature_html)
                    signature_images = []
                    for match in matches:
                        # Construire chemin complet vers l'image
                        img_path = os.path.join(signature_dir, match.replace('%20', ' ').replace('/', '\\'))
                        if os.path.exists(img_path):
                            cid = os.path.basename(match)  # Identifiant unique pour l'image dans le mail
                            # Modifier la source dans le HTML signature pour utiliser cid
                            signature_html = signature_html.replace(f'src="{match}"', f'src="cid:{cid}"')
                            signature_images.append((img_path, cid))
                        else:
                            logging.warning(f"Signature image not found: {img_path}")
                except Exception as e:
                    logging.error(f"Error reading signature:  {e}")


            # 2. Lire la signature HTML
            with open(signature_path, "r", encoding="mbcs") as f:
                signature_html = f.read()

            # 3. Chercher les images référencées dans la signature
            matches = re.findall(r'<img[^>]+src="([^"]+)"', signature_html)
            signature_images = []
            for match in matches:
                # Construire chemin complet vers l'image
                img_path = os.path.join(signature_dir, match.replace('%20', ' ').replace('/', '\\'))
                if os.path.exists(img_path):
                    cid = os.path.basename(match)  # Identifiant unique pour l'image dans le mail
                    # Modifier la source dans le html signature pour utiliser cid
                    signature_html = signature_html.replace(f'src="{match}"', f'src="cid:{cid}"')
                    signature_images.append((img_path, cid))
                else:
                    logging.warning(f"Signature image not found:  {img_path}")

            # 4. Créer un mail Outlook
            outlook = win32.Dispatch('outlook.application')
            mail = outlook.CreateItem(0)
            mail.To = destinataire
            mail.Subject = sujet
            if cc:
                mail.CC = cc

            # 5. Construire le corps HTML avec le message + la signature
            message_html = corps.replace('\n', '<br>')
            full_body = f"<html><body>{message_html}<br><br>{signature_html}</body></html>"
            mail.HTMLBody = full_body

            # 6. Ajouter les images en pièces jointes liées avec leur cid
            for img_path, cid in signature_images:
                attachment = mail.Attachments.Add(img_path)
                # Fixe la propriété MAPI pour que l'image s'affiche inline
                attachment.PropertyAccessor.SetProperty("http://schemas.microsoft.com/mapi/proptag/0x3712001F", cid)

            # 7. Ajouter un fichier joint si fourni
            if fichier_joint:
                mail.Attachments.Add(fichier_joint)

            # 8. Envoyer le mail
            mail.Send()
            return True
        except Exception as e:
            logging.error(f"Error sending email via Outlook:  {e}")
            return False
        


    def run(self):
        try:
            # Vérification existence du modèle
            if not os.path.exists(self.template_path):
                raise FileNotFoundError(f"PPTX template not found : {self.template_path}")
            
            df = pd.read_excel(self.excel_path, skiprows=1)
            date_edition = self.format_date_with_suffix(datetime.today())
            self.progress.emit("🕒 Generating certificates...")
            #total_certificates = 0


            if not os.path.exists(self.output_folder):
                os.makedirs(self.output_folder)
                logging.info(f"Output folder created: {self.output_folder}")

            total_rows = len(df)

            for index, row in df.iterrows():
                
                # === Gestion de la date de formation ===
                try:
                    formation_date = pd.to_datetime(str(row.iloc[0])).date()
                except Exception:
                    msg = f"Row {index + 2}: invalid training date, certificate ignored."
                    self.progress.emit(f"⚠️ {msg}")
                    logging.warning(msg)
                    continue

                # Filtrage par date
                if self.date_start and formation_date < self.date_start:
                    continue
                if self.date_end and formation_date > self.date_end:
                    continue

                # === Gestion du score ===
                try:
                    score = float(row.iloc[3])
                except Exception:
                    msg = f"Row {index + 2}: invalid score, certificate ignored."
                    self.progress.emit(f"⚠️ {msg}")
                    logging.warning(msg)
                    continue

                if score < self.score_min:
                    msg = f"{row.iloc[1]} not certified (score {score} < {self.score_min})"

                    logging.info(msg)
                    continue

                nom = str(row.iloc[1])

                # Gestion format date
                if isinstance(row.iloc[0], pd.Timestamp):
                    date_formation = row.iloc[0].strftime('%d/%m/%Y')
                else:
                    try:
                        date_formation = pd.to_datetime(str(row.iloc[0])).strftime('%d/%m/%Y')
                    except Exception:
                        date_formation = str(row.iloc[0])

                try:
                    prs = Presentation(self.template_path)
                except Exception as e:
                    raise ValueError(f"Cannot load template: {self.template_path} - Error: {e}")

                sso_raw = row.iloc[7]
                email_destinataire = str(row.iloc[8])  # Colonne email

                if pd.isna(sso_raw):
                    sso_utilisateur = ""
                else:
                    sso_utilisateur = str(int(sso_raw)) if isinstance(sso_raw, float) and sso_raw.is_integer() else str(sso_raw)


                # Remplacement des balises
                replacements = {
                    "{{NOM}}": nom,
                    "{{SSO}}": sso_utilisateur,
                    "{{FORMATION}}": self.formation_title,
                    "{{DATE_FORMATION}}": date_formation,
                    "{{DATE_EDITION}}": date_edition
                }
                self.replace_text(prs, replacements)

                # Construction des chemins
                safe_nom = self.clean_filename(nom)
                safe_title = self.clean_filename(self.formation_title)

                pptx_path = os.path.abspath(os.path.join(self.output_folder, f"{safe_title} - {safe_nom}.pptx"))
                pdf_path = os.path.abspath(os.path.join(self.output_folder, f"{safe_title} - {safe_nom}.pdf"))

                try:
                    prs.save(pptx_path)
                    logging.info(f"PPTX saved: {pptx_path}")
                except Exception as e:
                    self.progress.emit(f"❌ PPTX save error : {e}")
                    logging.error(f"PPTX save error for {nom} : {e}")
                    continue

                self.convert_to_pdf(pptx_path, pdf_path)

                if os.path.exists(pdf_path):
                    try:
                        os.remove(pptx_path)
                        logging.info(f"PPTX deleted : {pptx_path}")
                    except Exception as e:
                        warning_msg = f"Could not delete {pptx_path}: {e}"
                        self.progress.emit(f"⚠️ {warning_msg}")
                        logging.warning(warning_msg)
                    logging.info(f"Certificate generated for: {nom}")
                    self.total_certificates += 1
                else:
                    error_msg = f"PDF not generated for {nom}."
                    self.progress.emit(f"❌ {error_msg}")
                    logging.error(error_msg)
        
                copie_email = "SAM.L3SSystem@gevernova.com"
                if email_destinataire and os.path.exists(pdf_path):
                    mois_formation = datetime.today().strftime("%B %Y") 
                    sujet = f"{mois_formation}, {self.formation_title} - Training Certificate"
                    corps = f"""
                            Dears,

                            I am pleased to share with you your {self.formation_title}  certificate. Please find attached your certificate.

                            Best regards
                            """
                    
                    if self.envoyer_email_outlook(email_destinataire, sujet, corps, pdf_path, cc=copie_email):
                        self.progress.emit(f"✉️Certificate sent to {email_destinataire}")
                        logging.info(f"Certificate sent to  {email_destinataire} (cc {copie_email})")
                    else:
                        self.progress.emit(f"❌Failed to send email to {email_destinataire}")
                        logging.error(f"Failed to send email to{email_destinataire}")

                 # En fin de boucle, calcul de la progression
                percent = int((index + 10) / total_rows * 100)
                self.progress_percent.emit(percent)

            self.finished.emit(f"✅ {self.total_certificates} certificates successfully generated.")
            logging.info(f"{self.total_certificates} certificates generated successfully.")

        except Exception as e:
            error_msg = f"General error : {e}"
            self.finished.emit(f"❌ {error_msg}")
            logging.error(error_msg)
